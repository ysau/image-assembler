import os

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

locations = ['Loc_1', 'Loc_2', 'Loc_3', 'Loc_4', 'Loc_5', 'Loc_6', 'Loc_7', 'Special']

image_width = 1
image_height = 1
margin_between_images = 0.05

title_height = 0.9
first_row_height = 0.25
first_column_width = 1.5
data_row_height = 0.4

def generate_report(path='.', template='template.pptx'):
    prs = Presentation(template)

    samples = [f for f in os.listdir(path) if os.path.isdir(os.path.join(path, f))]
    for sample in samples:
        add_slide(prs, path, sample)

    prs.save('output.pptx')


def add_slide(prs, path, sample):
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    add_table(slide, sample)
    add_images(slide, path, sample)


def add_table(slide, sample):
    rows = 3
    cols = 9
    left = Inches(margin_between_images)
    top = Inches(title_height + margin_between_images)
    width = Inches(6.0)
    height = Inches(2.0)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    table.rows[0].height = Inches(first_row_height)
    table.rows[1].height = Inches(image_height + margin_between_images)
    table.rows[2].height = Inches(image_height + margin_between_images)

    # set column widths
    table.columns[0].width = Inches(first_column_width)
    for i in range(1, 9):
        table.columns[i].width = Inches(image_height + margin_between_images)

    # write column headings
    table.cell(0, 1).text = 'Loc A'
    for ind, label in enumerate(locations[1:-1]):
        table.cell(0, ind + 2).text = locations[1:-1][ind]
    table.cell(0, 8).text = 'Special'

    for cell in table.rows[0].cells:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

    table.cell(2, 0).text = sample
    for paragraph in table.rows[2].cells[0].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)


def add_images(slide, path, sample):
    for ind, loc in enumerate(locations):
        try:
            if loc == 'Special':
                loc = 'Loc_Special'
                try:
                    filename = [f for f in os.listdir(os.path.join(path, sample, loc)) if f[-4:] == '.tif'][0]
                except IndexError:
                    raise FileNotFoundError
            else:
                try:
                    filename = [f for f in os.listdir(os.path.join(path, sample, loc)) if f[-4:] == '.tif'][0]
                except IndexError:
                    raise FileNotFoundError
            image_path = os.path.join(path, sample, loc, filename)
 
            left = Inches(margin_between_images + first_column_width + margin_between_images/2 + ind * (image_width + margin_between_images))
            top = Inches(title_height + data_row_height + margin_between_images/2 + image_height + margin_between_images)
            height = width = Inches(1)
            slide.shapes.add_picture(image_path, left, top, height, width)
        except FileNotFoundError as not_found:
            print('Skipped:', not_found.filename)


def generate_ex_images():

    def generate_number_image(i):
        img = Image.new('RGB', (100, 100), color=(73, 109, 137))
        d = ImageDraw.Draw(img)
        d.text((10, 10), "Image {}".format(i), fill=(255, 255, 0))
        img.save(os.path.join('ex_images', 'sample1', 'Loc_{}'.format(i), 'number_{}.tif'.format(i)))

    os.makedirs('ex_images')
    os.makedirs(os.path.join('ex_images', 'sample1'))
    for i in range(1, 8):
        os.makedirs(os.path.join('ex_images', 'sample1', 'Loc_{}'.format(i)))
        generate_number_image(i)
    os.makedirs(os.path.join('ex_images', 'sample1', 'Loc_Special'))
    generate_number_image('Special')


if __name__ == '__main__':
    generate_ex_images()
    generate_report('ex_images')